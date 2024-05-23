import os
import tkinter as tk
from tkinter import filedialog, simpledialog
from pptx import Presentation
from docx import Document
import openpyxl
import fitz
from googletrans import Translator
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
class FileTranslator:
   def __init__(self):
       self.root = tk.Tk()
       self.root.withdraw()  # Hide the root window
   def select_input_file(self):
       # Open file dialog to select input file
       file_path = filedialog.askopenfilename(title="Select Input File")
       return file_path
   def specify_target_language(self):
       # Prompt user to enter target language code
       target_language = simpledialog.askstring("Target Language", "Enter the target language code (e.g., 'es' for Spanish): ")
       if target_language:
           target_language = target_language.strip()
       return target_language
   def extract_text_from_docx(self, file_path):
       # Extract text from DOCX file
       if not os.path.exists(file_path):
           print(f"Error: Input file '{file_path}' not found.")
           return None
       doc = Document(file_path)
       extracted_text = ""
       for paragraph in doc.paragraphs:
           extracted_text += paragraph.text + "\n"
       return extracted_text
   def extract_text_from_pptx(self, file_path):
       # Extract text from PPTX file
       if not os.path.exists(file_path):
           print(f"Error: Input file '{file_path}' not found.")
           return None
       presentation = Presentation(file_path)
       extracted_text = ""
       for slide in presentation.slides:
           for shape in slide.shapes:
               if hasattr(shape, "text"):
                   extracted_text += shape.text + "\n"
       return extracted_text
   def extract_text_from_pdf(self, file_path):
       # Extract text from PDF file
       if not os.path.exists(file_path):
           print(f"Error: Input file '{file_path}' not found.")
           return None
       extracted_text = ""
       with fitz.open(file_path) as pdf_document:
           for page in pdf_document:
               extracted_text += page.get_text()
       return extracted_text
   def extract_text_from_xlsx(self, file_path):
       # Extract text from XLSX file
       if not os.path.exists(file_path):
           print(f"Error: Input file '{file_path}' not found.")
           return ""
       extracted_text = ""
       try:
           wb = openpyxl.load_workbook(file_path)
           for sheet in wb.sheetnames:
               ws = wb[sheet]
               for row in ws.iter_rows(values_only=True):
                   for cell in row:
                       if cell:
                           extracted_text += str(cell) + "\n"
       except Exception as e:
           print(f"Error extracting text from XLSX: {e}")
       return extracted_text
   def extract_text_from_txt(self, file_path):
       # Extract text from TXT file
       if not os.path.exists(file_path):
           print(f"Error: Input file '{file_path}' not found.")
           return None
       extracted_text = ""
       with open(file_path, 'r') as file:
           extracted_text = file.read()
       return extracted_text
   def extract_text_from_properties(self, file_path):
       # Extract text from properties file
       if not os.path.exists(file_path):
           print(f"Error: Input file '{file_path}' not found.")
           return None
       extracted_text = ""
       with open(file_path, 'r') as file:
           for line in file:
               line = line.strip()
               if line and not line.startswith('#'):  # Skip comments and blank lines
                   key_value = line.split('=', 1)
                   if len(key_value) == 2:
                       key, value = key_value
                       extracted_text += f"{key.strip()}={value.strip()}\n"
       return extracted_text
   def translate_text(self, text, target_language):
       # Translate text to target language using Google Translate
       try:
           translator = Translator()
           translated = translator.translate(text, dest=target_language)
           return translated.text
       except Exception as e:
           print(f"Translation failed: {e}")
           return None
   def save_translated_text(self, translated_text, file_path):
       # Construct translated file path
       file_name, file_extension = os.path.splitext(os.path.basename(file_path))
       translated_file_path = os.path.join(os.path.dirname(file_path), f"{file_name}_translated{file_extension}")
       try:
           if file_extension.lower() == '.docx':
               doc = Document()
               for paragraph in translated_text.split('\n'):
                   doc.add_paragraph(paragraph)
               doc.save(translated_file_path)
           elif file_extension.lower() == '.pptx':
               presentation = Presentation()
               for slide_text in translated_text.split('\n\n'):
                   slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                   text_frame = slide.shapes.add_textbox(0, 0, 1, 1).text_frame
                   text_frame.text = slide_text
               presentation.save(translated_file_path)
           elif file_extension.lower() == '.pdf':
               # Create a new PDF file
               styles = getSampleStyleSheet()
               doc = SimpleDocTemplate(translated_file_path, pagesize=letter)
               elements = []
               for paragraph in translated_text.split('\n'):
                   elements.append(Paragraph(paragraph, styles["BodyText"]))
               doc.build(elements)
           elif file_extension.lower() == '.xlsx':
               # Write translated text to a new Excel file
               translated_workbook = openpyxl.Workbook()
               translated_sheet = translated_workbook.active
               for line in translated_text.split('\n'):
                   translated_sheet.append([line])
               translated_workbook.save(translated_file_path)
           elif file_extension.lower() == '.txt':
               # Write translated text to a new TXT file
               with open(translated_file_path, 'w') as file:
                   file.write(translated_text)
           else:
               print("Unsupported file format.")
               return None
       except Exception as e:
           print(f"Error saving translated file: {e}")
           return None
       return translated_file_path
   def run(self):
       # Select input file
       input_file = self.select_input_file()
       if not input_file:
           print("No input file selected.")
           return
       # Specify target language
       target_language = self.specify_target_language()
       # Extract text based on file extension
       file_extension = os.path.splitext(input_file)[-1].lower()
       if file_extension == '.docx':
           extracted_text = self.extract_text_from_docx(input_file)
       elif file_extension == '.pptx':
           extracted_text = self.extract_text_from_pptx(input_file)
       elif file_extension == '.pdf':
           extracted_text = self.extract_text_from_pdf(input_file)
       elif file_extension == '.xlsx':
           extracted_text = self.extract_text_from_xlsx(input_file)
       elif file_extension == '.txt':
           extracted_text = self.extract_text_from_txt(input_file)
       elif file_extension == '.properties':
           extracted_text = self.extract_text_from_properties(input_file)
       else:
           print("Unsupported file format.")
           return
       if extracted_text:
           # Translate extracted text
           translated_text = self.translate_text(extracted_text, target_language)
           if translated_text:
               # Save translated text to a file
               translated_file_path = self.save_translated_text(translated_text, input_file)
               if translated_file_path:
                   print(f"Translated content saved as: {translated_file_path}")
               else:
                   print("Failed to save translated content.")
           else:
               print("Translation failed.")
       else:
           print("Failed to extract text from input file.")
# Main program
if __name__ == "__main__":
   translator = FileTranslator()
   translator.run()